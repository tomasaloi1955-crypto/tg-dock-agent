"""
TG Doc Agent v3 — красивые документы всех форматов со стилями.
PDF, DOCX, XLSX, CSV, TXT, PPTX — у каждого 4 стиля.
ИИ: Groq → Gemini fallback.
"""

import io
import json
import logging
import os

import requests
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    Application, CommandHandler, MessageHandler,
    CallbackQueryHandler, ContextTypes, filters,
)

logging.basicConfig(format="%(asctime)s %(levelname)s %(message)s", level=logging.INFO)
log = logging.getLogger("doc-agent")
logging.getLogger("httpx").setLevel(logging.WARNING)
logging.getLogger("httpcore").setLevel(logging.WARNING)

BOT_TOKEN = os.environ["BOT_TOKEN"]
GEMINI_KEY = os.environ.get("GEMINI_API_KEY", "")
GEMINI_MODEL = os.environ.get("GEMINI_MODEL", "gemini-2.0-flash")
GROQ_KEY = os.environ.get("GROQ_API_KEY", "")
GROQ_MODEL = os.environ.get("GROQ_MODEL", "llama-3.3-70b-versatile")
WEBHOOK_URL = os.environ.get("WEBHOOK_URL", "")
PORT = int(os.environ.get("PORT", "10000"))

TEMPLATES_FILE = "templates.json"
FONT_PATH = os.path.join(os.path.dirname(__file__), "fonts", "DejaVuSans.ttf")
FONT_BOLD_PATH = os.path.join(os.path.dirname(__file__), "fonts", "DejaVuSans-Bold.ttf")

# ── Палитры стилей ────────────────────────────────────────────────────────
STYLES = {
    "minimal": {
        "name": "🌊 Минимализм",
        "bg": "F8F9FA", "header_bg": "212529", "header_fg": "FFFFFF",
        "accent": "4ECDC4", "text": "212529", "subtext": "6C757D",
        "row_alt": "F1F3F5",
    },
    "business": {
        "name": "💎 Бизнес",
        "bg": "1E2761", "header_bg": "0D1B4E", "header_fg": "CADCFC",
        "accent": "F5A623", "text": "CADCFC", "subtext": "8FA8D8",
        "row_alt": "253180",
    },
    "creative": {
        "name": "🎨 Яркий",
        "bg": "FFFFFF", "header_bg": "F96167", "header_fg": "FFFFFF",
        "accent": "2F3C7E", "text": "2F3C7E", "subtext": "6B7280",
        "row_alt": "FFF0F0",
    },
    "edu": {
        "name": "🌍 Образование",
        "bg": "FFFFFF", "header_bg": "028090", "header_fg": "FFFFFF",
        "accent": "02C39A", "text": "2C3E50", "subtext": "7F8C8D",
        "row_alt": "E8F8F5",
    },
}


def hex2rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


# ── Шаблоны ───────────────────────────────────────────────────────────────
def load_templates():
    if os.path.exists(TEMPLATES_FILE):
        with open(TEMPLATES_FILE, encoding="utf-8") as f:
            return json.load(f)
    return {}


def save_template(user_id, text):
    data = load_templates()
    data[str(user_id)] = text[:8000]
    with open(TEMPLATES_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False)


def get_template(user_id):
    return load_templates().get(str(user_id), "")


# ── LLM ───────────────────────────────────────────────────────────────────
def call_groq(prompt, system):
    r = requests.post(
        "https://api.groq.com/openai/v1/chat/completions",
        headers={"Authorization": f"Bearer {GROQ_KEY}"},
        json={"model": GROQ_MODEL, "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ], "temperature": 0.7},
        timeout=120,
    )
    if r.status_code != 200:
        raise RuntimeError(f"Groq {r.status_code}")
    return r.json()["choices"][0]["message"]["content"]


def call_gemini(prompt, system):
    import time
    body = {"contents": [{"parts": [{"text": prompt}]}]}
    if system:
        body["systemInstruction"] = {"parts": [{"text": system}]}
    for model in [GEMINI_MODEL, "gemini-2.0-flash-lite"]:
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
        for _ in range(2):
            r = requests.post(url, json=body, headers={"x-goog-api-key": GEMINI_KEY}, timeout=120)
            if r.status_code == 429:
                time.sleep(5); continue
            if r.status_code == 200:
                return r.json()["candidates"][0]["content"]["parts"][0]["text"]
            break
    raise RuntimeError("Gemini недоступен")


def llm(prompt, system=""):
    if GROQ_KEY:
        try: return call_groq(prompt, system)
        except Exception as e: log.warning("Groq: %s", e)
    if GEMINI_KEY:
        try: return call_gemini(prompt, system)
        except Exception as e: log.warning("Gemini: %s", e)
    raise RuntimeError("Нет доступных ИИ-провайдеров")


PLAN_SYSTEM = """Ты — генератор документов. Верни ТОЛЬКО валидный JSON без markdown-обёрток:
{
 "format": "pdf"|"docx"|"xlsx"|"csv"|"txt"|"pptx",
 "filename": "короткое_имя",
 "title": "Заголовок",
 "content": ...
}
Формат: презентация/слайды→pptx, таблица/данные→xlsx/csv, гайд/чек-лист/инструкция→pdf, ворд→docx, текст→txt
Content:
- pdf/docx/txt: [{"type":"heading","text":"..."},{"type":"paragraph","text":"..."},{"type":"bullets","items":["..."]},{"type":"checklist","items":["..."]}]
- xlsx/csv: {"headers":["..."],"rows":[["..."]]}
- pptx: [{"title":"...","bullets":["...","..."]}]
Делай контент содержательным и полным. Для pptx — минимум 5 слайдов.
Если есть ОБРАЗЕЦ — копируй его структуру и стиль."""


def plan_document(user_request, template=""):
    prompt = user_request
    if template:
        prompt += f"\n\n--- ОБРАЗЕЦ ---\n{template}"
    raw = llm(prompt, PLAN_SYSTEM).strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        raw = raw[4:] if raw.startswith("json") else raw
    return json.loads(raw.strip())


# ── PDF со стилями ────────────────────────────────────────────────────────
def make_pdf(plan, style_key="minimal"):
    from fpdf import FPDF
    s = STYLES[style_key]
    bg = hex2rgb(s["bg"])
    hbg = hex2rgb(s["header_bg"])
    hfg = hex2rgb(s["header_fg"])
    acc = hex2rgb(s["accent"])
    txt = hex2rgb(s["text"])
    sub = hex2rgb(s["subtext"])

    pdf = FPDF(format="A4")
    pdf.add_font("DejaVu", "", FONT_PATH)
    pdf.add_font("DejaVu", "B", FONT_BOLD_PATH)
    pdf.set_auto_page_break(auto=True, margin=22)
    pdf.add_page()

    # Шапка — цветной блок на всю ширину
    pdf.set_fill_color(*hbg)
    pdf.rect(0, 0, 210, 28, "F")
    pdf.set_font("DejaVu", "B", 18)
    pdf.set_text_color(*hfg)
    pdf.set_xy(10, 6)
    pdf.cell(190, 16, plan["title"], align="L")

    # Цветная линия-разделитель под шапкой
    pdf.set_fill_color(*acc)
    pdf.rect(0, 28, 210, 3, "F")
    pdf.ln(10)

    for block in plan["content"]:
        t = block["type"]
        if t == "heading":
            pdf.ln(3)
            # Акцентная полоска слева
            pdf.set_fill_color(*acc)
            pdf.rect(10, pdf.get_y(), 4, 9, "F")
            pdf.set_font("DejaVu", "B", 13)
            pdf.set_text_color(*txt)
            pdf.set_x(17)
            pdf.multi_cell(180, 9, block["text"], new_x="LMARGIN", new_y="NEXT")
            pdf.ln(1)

        elif t == "paragraph":
            pdf.set_font("DejaVu", "", 11)
            pdf.set_text_color(*txt)
            pdf.set_x(10)
            pdf.multi_cell(190, 6, block["text"], new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)

        elif t == "bullets":
            pdf.set_font("DejaVu", "", 11)
            pdf.set_text_color(*txt)
            for item in block["items"]:
                # Цветная точка
                pdf.set_fill_color(*acc)
                pdf.rect(12, pdf.get_y() + 2, 3, 3, "F")
                pdf.set_x(18)
                pdf.multi_cell(182, 6, item, new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)

        elif t == "checklist":
            pdf.set_font("DejaVu", "", 11)
            pdf.set_text_color(*txt)
            for item in block["items"]:
                # Квадратик-чекбокс
                pdf.set_draw_color(*acc)
                pdf.set_line_width(0.5)
                pdf.rect(12, pdf.get_y() + 1, 4, 4)
                pdf.set_x(19)
                pdf.multi_cell(181, 6, item, new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)

    # Подвал
    pdf.set_y(-15)
    pdf.set_font("DejaVu", "", 8)
    pdf.set_text_color(*sub)
    pdf.set_fill_color(*hbg)
    pdf.rect(0, pdf.get_y() - 2, 210, 20, "F")
    pdf.set_text_color(*hfg)
    pdf.cell(0, 8, plan["title"], align="C")

    return bytes(pdf.output())


# ── DOCX со стилями ───────────────────────────────────────────────────────
def make_docx(plan, style_key="minimal"):
    from docx import Document
    from docx.shared import Pt, RGBColor as DRGBColor, Inches, Cm
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    import lxml.etree as etree

    s = STYLES[style_key]
    acc = DRGBColor(*hex2rgb(s["accent"]))
    hbg = DRGBColor(*hex2rgb(s["header_bg"]))
    hfg = DRGBColor(*hex2rgb(s["header_fg"]))
    txt = DRGBColor(*hex2rgb(s["text"]))

    doc = Document()
    # Поля
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # Заголовок документа
    title_p = doc.add_paragraph()
    title_p.paragraph_format.space_after = Pt(6)
    run = title_p.add_run(plan["title"])
    run.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = hfg
    # Цвет фона параграфа заголовка через XML
    pPr = title_p._p.get_or_add_pPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), s["header_bg"])
    pPr.append(shd)

    doc.add_paragraph()  # отступ

    for block in plan["content"]:
        t = block["type"]
        if t == "heading":
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(10)
            p.paragraph_format.space_after = Pt(4)
            # Левая граница — акцентный цвет
            pPr = p._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            left = OxmlElement("w:left")
            left.set(qn("w:val"), "single")
            left.set(qn("w:sz"), "24")
            left.set(qn("w:space"), "8")
            left.set(qn("w:color"), s["accent"])
            pBdr.append(left)
            pPr.append(pBdr)
            run = p.add_run(block["text"])
            run.bold = True
            run.font.size = Pt(13)
            run.font.color.rgb = acc

        elif t == "paragraph":
            p = doc.add_paragraph(block["text"])
            p.paragraph_format.space_after = Pt(6)
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = txt

        elif t == "bullets":
            for item in block["items"]:
                p = doc.add_paragraph(style="List Bullet")
                run = p.add_run(item)
                run.font.size = Pt(11)
                run.font.color.rgb = txt

        elif t == "checklist":
            for item in block["items"]:
                p = doc.add_paragraph()
                # Цветной чекбокс
                cb = p.add_run("☐  ")
                cb.font.color.rgb = acc
                cb.font.size = Pt(12)
                run = p.add_run(item)
                run.font.size = Pt(11)
                run.font.color.rgb = txt

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── XLSX со стилями ───────────────────────────────────────────────────────
def make_xlsx(plan, style_key="minimal"):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, GradientFill
    from openpyxl.styles.numbers import FORMAT_NUMBER_COMMA_SEPARATED1

    s = STYLES[style_key]
    wb = Workbook()
    ws = wb.active
    ws.title = plan["title"][:30]

    c = plan["content"]
    headers = c["headers"]
    rows = c["rows"]

    # Заголовок листа
    ws.merge_cells(f"A1:{chr(64+len(headers))}1")
    title_cell = ws["A1"]
    title_cell.value = plan["title"]
    title_cell.font = Font(bold=True, size=14, color=s["header_fg"])
    title_cell.fill = PatternFill("solid", fgColor=s["header_bg"])
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 30

    # Шапка таблицы
    thin = Side(style="thin", color=s["accent"])
    border = Border(bottom=thin)
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=2, column=col, value=h)
        cell.font = Font(bold=True, color=s["header_fg"])
        cell.fill = PatternFill("solid", fgColor=s["accent"])
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = border
    ws.row_dimensions[2].height = 22

    # Данные с чередующимися строками
    for r_idx, row in enumerate(rows, 3):
        alt = r_idx % 2 == 0
        fill_color = s["row_alt"] if alt else s["bg"]
        for c_idx, val in enumerate(row, 1):
            cell = ws.cell(row=r_idx, column=c_idx, value=val)
            cell.fill = PatternFill("solid", fgColor=fill_color)
            cell.font = Font(color=s["text"])
            cell.alignment = Alignment(vertical="center", wrap_text=True)
        ws.row_dimensions[r_idx].height = 18

    # Ширина столбцов
    for col in ws.columns:
        if col[0].row == 1:
            continue
        width = max(len(str(cell.value or "")) for cell in col) + 4
        ws.column_dimensions[col[0].column_letter].width = min(width, 45)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── CSV ───────────────────────────────────────────────────────────────────
def make_csv(plan, style_key=None):
    import csv
    c = plan["content"]
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow(c["headers"])
    w.writerows(c["rows"])
    return buf.getvalue().encode("utf-8-sig")


# ── TXT со стилями ────────────────────────────────────────────────────────
def make_txt(plan, style_key="minimal"):
    borders = {
        "minimal": ("─", "│", "┌", "┐", "└", "┘"),
        "business": ("═", "║", "╔", "╗", "╚", "╝"),
        "creative": ("*", "|", "+", "+", "+", "+"),
        "edu": ("-", "|", "+", "+", "+", "+"),
    }
    h, v, tl, tr, bl, br = borders.get(style_key, borders["minimal"])
    w = 60
    lines = []

    # Рамка заголовка
    lines.append(tl + h * (w - 2) + tr)
    title = plan["title"]
    pad = (w - 2 - len(title)) // 2
    lines.append(v + " " * pad + title + " " * (w - 2 - pad - len(title)) + v)
    lines.append(bl + h * (w - 2) + br)
    lines.append("")

    for block in plan["content"]:
        t = block["type"]
        if t == "heading":
            lines += ["", f"  {block['text'].upper()}", f"  {'─' * len(block['text'])}", ""]
        elif t == "paragraph":
            lines += [f"  {block['text']}", ""]
        elif t == "bullets":
            for item in block["items"]:
                lines.append(f"  ► {item}")
            lines.append("")
        elif t == "checklist":
            for item in block["items"]:
                lines.append(f"  [ ] {item}")
            lines.append("")

    lines += ["", tl + h * (w - 2) + tr,
              v + " Создано TG Doc Agent".center(w - 2) + v,
              bl + h * (w - 2) + br]
    return "\n".join(lines).encode("utf-8")


# ── PPTX со стилями ───────────────────────────────────────────────────────
def make_pptx(plan, style_key="minimal"):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    s = STYLES[style_key]
    bg_rgb = RGBColor(*hex2rgb(s["bg"]))
    hbg_rgb = RGBColor(*hex2rgb(s["header_bg"]))
    hfg_rgb = RGBColor(*hex2rgb(s["header_fg"]))
    acc_rgb = RGBColor(*hex2rgb(s["accent"]))
    txt_rgb = RGBColor(*hex2rgb(s["text"]))

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    def rect(slide, x, y, w, h, color, line=False):
        shape = slide.shapes.add_shape(1, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        if not line:
            shape.line.fill.background()
        return shape

    def textbox(slide, x, y, w, h, text, bold=False, size=18, color=None, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(size)
        if color:
            run.font.color.rgb = color
        return tb

    slides_data = plan["content"]

    # Титульный слайд
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hbg_rgb
    rect(slide, 0, int(H * 0.72), W, int(H * 0.28), acc_rgb)
    textbox(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.5),
            plan["title"], bold=True, size=44, color=hfg_rgb)
    if slides_data and slides_data[0].get("bullets"):
        textbox(slide, Inches(0.8), Inches(4.3), Inches(10), Inches(0.9),
                slides_data[0]["bullets"][0], size=20, color=acc_rgb)

    # Контентные слайды
    for i, s_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_rgb

        # Шапка
        rect(slide, 0, 0, W, Inches(1.3), hbg_rgb)
        rect(slide, 0, 0, Inches(0.55), Inches(1.3), acc_rgb)
        textbox(slide, Inches(0.05), Inches(0.25), Inches(0.45), Inches(0.8),
                str(i + 1), bold=True, size=20, color=hfg_rgb, align=PP_ALIGN.CENTER)
        textbox(slide, Inches(0.75), Inches(0.15), Inches(11.6), Inches(1.05),
                s_data.get("title", ""), bold=True, size=26, color=hfg_rgb)

        # Буллиты
        bullets = s_data.get("bullets", [])
        y0 = Inches(1.55)
        h_each = Inches(0.9) if len(bullets) > 4 else Inches(1.05)
        for j, b in enumerate(bullets[:6]):
            y = y0 + j * h_each
            rect(slide, Inches(0.48), y + Inches(0.2), Inches(0.16), Inches(0.16), acc_rgb)
            textbox(slide, Inches(0.8), y, Inches(11.5), h_each, b, size=17, color=txt_rgb)

        # Декор справа
        rect(slide, W - Inches(0.22), Inches(1.3), Inches(0.22), H - Inches(1.3), acc_rgb)

    # Финальный слайд
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hbg_rgb
    rect(slide, 0, 0, W, Inches(0.45), acc_rgb)
    rect(slide, 0, H - Inches(0.45), W, Inches(0.45), acc_rgb)
    textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(2),
            "Спасибо!", bold=True, size=54, color=hfg_rgb, align=PP_ALIGN.CENTER)
    textbox(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.9),
            plan["title"], size=20, color=acc_rgb, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


GENERATORS = {
    "pdf": make_pdf, "docx": make_docx, "xlsx": make_xlsx,
    "csv": make_csv, "txt": make_txt, "pptx": make_pptx,
}

# ── UI ────────────────────────────────────────────────────────────────────
PENDING = {}  # {user_id: {"request": str, "fmt": str}}

START_TEXT = (
    "Привет! Я создаю красивые документы 📄✨\n\n"
    "Напиши запрос, например:\n"
    "• «чек-лист запуска ТГ-канала»\n"
    "• «гайд по prompt-инжинирингу»\n"
    "• «презентация про нейросети, 5 слайдов»\n"
    "• «таблица учёта расходов»\n\n"
    "📎 Пришли файл-образец (.txt, .docx, .pdf) — буду копировать стиль.\n\n"
    "/template — текущий образец\n/clear — удалить образец"
)

FORMAT_KB = InlineKeyboardMarkup([
    [InlineKeyboardButton("📊 Презентация", callback_data="fmt_pptx"),
     InlineKeyboardButton("📄 PDF", callback_data="fmt_pdf")],
    [InlineKeyboardButton("📝 Word", callback_data="fmt_docx"),
     InlineKeyboardButton("📊 Excel", callback_data="fmt_xlsx")],
    [InlineKeyboardButton("📋 CSV", callback_data="fmt_csv"),
     InlineKeyboardButton("📃 TXT", callback_data="fmt_txt")],
])

STYLE_KB = InlineKeyboardMarkup([
    [InlineKeyboardButton("🌊 Минимализм", callback_data="sty_minimal"),
     InlineKeyboardButton("💎 Бизнес", callback_data="sty_business")],
    [InlineKeyboardButton("🎨 Яркий", callback_data="sty_creative"),
     InlineKeyboardButton("🌍 Образование", callback_data="sty_edu")],
])

# CSV не имеет стилей — сразу генерируем
NO_STYLE_FMTS = {"csv"}


# ── Handlers ──────────────────────────────────────────────────────────────
async def cmd_start(u: Update, c): await u.message.reply_text(START_TEXT)


async def cmd_template(u: Update, c):
    t = get_template(u.effective_user.id)
    await u.message.reply_text(
        f"Образец (первые 500 симв.):\n\n{t[:500]}" if t else "Образец не загружен. Пришли файл!"
    )


async def cmd_clear(u: Update, c):
    d = load_templates()
    d.pop(str(u.effective_user.id), None)
    with open(TEMPLATES_FILE, "w", encoding="utf-8") as f:
        json.dump(d, f, ensure_ascii=False)
    await u.message.reply_text("Образец удалён ✅")


def extract_text(filename, data):
    name = filename.lower()
    if name.endswith((".txt", ".md", ".csv")):
        return data.decode("utf-8", errors="ignore")
    if name.endswith(".docx"):
        from docx import Document
        return "\n".join(p.text for p in Document(io.BytesIO(data)).paragraphs)
    if name.endswith(".pdf"):
        try:
            from pypdf import PdfReader
            return "\n".join(p.extract_text() or "" for p in PdfReader(io.BytesIO(data)).pages)
        except: return ""
    return ""


async def handle_document(u: Update, c):
    doc = u.message.document
    if doc.file_size > 5 * 1024 * 1024:
        await u.message.reply_text("Файл слишком большой (макс 5 МБ).")
        return
    f = await doc.get_file()
    data = bytes(await f.download_as_bytearray())
    text = extract_text(doc.file_name or "", data)
    if not text.strip():
        await u.message.reply_text("Не смог прочитать 😕 Жду: .txt, .md, .docx, .pdf, .csv")
        return
    save_template(u.effective_user.id, text)
    await u.message.reply_text("Образец сохранён ✅ Напиши что создать!")


async def handle_text(u: Update, c):
    uid = u.effective_user.id
    PENDING[uid] = {"request": u.message.text, "fmt": None}
    await u.message.reply_text(
        f"📋 *{u.message.text[:80]}*\n\nВыбери формат:",
        parse_mode="Markdown",
        reply_markup=FORMAT_KB,
    )


async def handle_callback(u: Update, c):
    q = u.callback_query
    await q.answer()
    uid = q.from_user.id
    data = q.data

    if data.startswith("fmt_"):
        fmt = data[4:]
        p = PENDING.get(uid)
        if not p:
            await q.edit_message_text("Сессия истекла. Напиши запрос заново."); return
        p["fmt"] = fmt

        if fmt in NO_STYLE_FMTS:
            await q.edit_message_text(f"⏳ Генерирую {fmt.upper()}...")
            await _gen(q, uid, fmt, None)
        else:
            fmt_names = {"pptx": "Презентацию", "pdf": "PDF", "docx": "Word",
                         "xlsx": "Excel", "txt": "TXT"}
            await q.edit_message_text(
                f"✅ Формат: {fmt_names.get(fmt, fmt)}\n\nВыбери стиль:",
                reply_markup=STYLE_KB,
            )

    elif data.startswith("sty_"):
        style_key = data[4:]
        p = PENDING.get(uid)
        if not p:
            await q.edit_message_text("Сессия истекла. Напиши запрос заново."); return
        fmt = p.get("fmt", "pdf")
        style_name = STYLES[style_key]["name"]
        await q.edit_message_text(f"⏳ Генерирую {fmt.upper()} в стиле {style_name}...")
        await _gen(q, uid, fmt, style_key)


async def _gen(q, uid, fmt, style_key):
    p = PENDING.get(uid, {})
    try:
        plan = plan_document(p.get("request", ""), get_template(uid))
        plan["format"] = fmt
        file_bytes = GENERATORS[fmt](plan, style_key or "minimal")
        filename = f"{plan.get('filename', 'document')}.{fmt}"
        style_label = f" [{STYLES[style_key]['name']}]" if style_key else ""
        await q.message.reply_document(
            document=io.BytesIO(file_bytes),
            filename=filename,
            caption=f"✅ {plan.get('title', 'Готово')}{style_label}",
        )
        PENDING.pop(uid, None)
    except Exception as e:
        log.exception("gen failed")
        await q.message.reply_text(f"Ошибка 😞 Попробуй переформулировать.\n({str(e)[:200]})")


def main():
    app = Application.builder().token(BOT_TOKEN).build()
    app.add_handler(CommandHandler("start", cmd_start))
    app.add_handler(CommandHandler("template", cmd_template))
    app.add_handler(CommandHandler("clear", cmd_clear))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))
    app.add_handler(CallbackQueryHandler(handle_callback))

    if WEBHOOK_URL:
        log.info("Webhook mode port %s", PORT)
        app.run_webhook(listen="0.0.0.0", port=PORT,
                        url_path=BOT_TOKEN,
                        webhook_url=f"{WEBHOOK_URL}/{BOT_TOKEN}")
    else:
        log.info("Polling mode")
        app.run_polling(drop_pending_updates=True)


if __name__ == "__main__":
    main()
